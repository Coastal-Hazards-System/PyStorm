"""Shared PyStorm slide-deck renderer: a Marp markdown deck -> editable native PPTX.

This is the standard pptx format for every PyStorm module. Each module keeps a local
(gitignored) Marp deck <Module>_Slides.md alongside its whitepaper and renders it with
this tool; the rendered .pptx is also local. Requires python-pptx (no LibreOffice or
browser), and produces real editable text boxes, tables, and pictures.

Usage (from a module directory):
    python ../../docs/slides_to_pptx.py <Module>_Slides.md <Module>_Slides.pptx

Layout rules that keep slides consistent:
  - one title bar + accent rule + footer on every content slide;
  - pure-text slides put all body blocks in ONE content frame spanning a fixed
    region (so paragraphs flow and never overlap);
  - image / table slides place those elements with measured heights;
  - shared typography (sizes, colors, spacing, native bullets).
"""
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn

SRC, OUT = sys.argv[1], sys.argv[2]

# ── Theme ────────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
INK = RGBColor(0x2B, 0x2B, 0x2B)
MUTE = RGBColor(0x80, 0x80, 0x80)
CODEINK = RGBColor(0x16, 0x4A, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODYF, CODEF = "Calibri", "Consolas"

MX = 0.75                      # left/right margin (in)
CW = 13.333 - 2 * MX           # content width
TITLE_Y, TITLE_H = 0.45, 0.80
RULE_Y = 1.28
BODY_Y, BODY_BOTTOM = 1.52, 6.95
FOOT_Y = 7.04

# ── Parse markdown into slides and blocks ────────────────────────────────────
text = open(SRC, encoding="utf-8").read()
if text.startswith("---"):
    text = text[text.index("\n---", 3) + 4:]
text = re.sub(r"<!--.*?-->", "", text, flags=re.S)
slides = [s.strip("\n") for s in re.split(r"\n---\n", text) if s.strip()]


def inline(s):
    s = re.sub(r"\*\*(.+?)\*\*", r"\1", s)
    s = re.sub(r"`(.+?)`", r"\1", s)
    s = re.sub(r"(?<!\*)\*(?!\s)(.+?)\*", r"\1", s)
    return s.strip()


def parse(lines):
    out, i, n = [], 0, len(lines)
    while i < n:
        ln = lines[i]
        if ln.startswith("```"):
            j, code = i + 1, []
            while j < n and not lines[j].startswith("```"):
                code.append(lines[j]); j += 1
            out.append(("code", code)); i = j + 1; continue
        m = re.match(r"!\[[^\]]*\]\(([^)]+)\)", ln)
        if m:
            out.append(("image", m.group(1))); i += 1; continue
        if ln.startswith("### "):
            out.append(("h3", inline(ln[4:]))); i += 1; continue
        if ln.startswith("## "):
            out.append(("title", inline(ln[3:]))); i += 1; continue
        if ln.startswith("# "):
            out.append(("big", inline(ln[2:]))); i += 1; continue
        if ln.lstrip().startswith("|") and ln.count("|") >= 2:
            rows = []
            while i < n and lines[i].lstrip().startswith("|"):
                rows.append(lines[i]); i += 1
            out.append(("table", rows)); continue
        if re.match(r"\s*[-*] ", ln):
            bs = []
            while i < n and re.match(r"\s*[-*] ", lines[i]):
                lead = len(lines[i]) - len(lines[i].lstrip())
                bs.append((min(lead // 2, 1), inline(re.sub(r"^\s*[-*] ", "", lines[i]))))
                i += 1
            out.append(("bullets", bs)); continue
        if re.match(r"\s*\d+\. ", ln):
            items = []
            while i < n and re.match(r"\s*\d+\. ", lines[i]):
                lead = len(lines[i]) - len(lines[i].lstrip())
                items.append((min(lead // 2, 1), inline(re.sub(r"^\s*\d+\.\s+", "", lines[i]))))
                i += 1
            out.append(("olist", items)); continue
        if not ln.strip():
            i += 1; continue
        para = []
        while i < n and lines[i].strip() and not re.match(r"\s*([-*] |\d+\. |#|```|\||!\[)", lines[i]):
            para.append(inline(lines[i])); i += 1
        out.append(("para", " ".join(para)))
    return out


# ── Build the presentation ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def style_run(p, s, size, bold=False, color=INK, mono=False, italic=False):
    r = p.add_run(); r.text = s; f = r.font
    f.size, f.bold, f.italic = Pt(size), bold, italic
    f.color.rgb = color; f.name = CODEF if mono else BODYF
    return r


def _bullet_pPr(p, level, indent=228600):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(274320 * (level + 1)))
    pPr.set("indent", str(-indent))
    for tag in ("a:buFont", "a:buChar", "a:buNone", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    return pPr


def bullet(p, level):
    pPr = _bullet_pPr(p, level)
    bf = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    bc = pPr.makeelement(qn("a:buChar"), {"char": "•" if level == 0 else "–"})
    pPr.append(bf); pPr.append(bc)


def numbered(p, level):
    pPr = _bullet_pPr(p, level, indent=320040)
    bf = pPr.makeelement(qn("a:buFont"), {"typeface": "+mj-lt"})
    ba = pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})
    pPr.append(bf); pPr.append(ba)


def add_blocks(tf, body, first=True):
    """Render text/code/h3/bullet blocks as paragraphs of one text frame."""
    for kind, payload in body:
        if kind == "h3":
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.space_before, p.space_after = Pt(12), Pt(5)
            style_run(p, payload, 22, True, NAVY)
        elif kind == "para":
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            p.space_after = Pt(9); p.line_spacing = 1.05
            style_run(p, payload, 19, color=INK)
        elif kind in ("bullets", "olist"):
            for lvl, txt in payload:
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                p.level = lvl
                (numbered if kind == "olist" else bullet)(p, lvl)
                p.space_after = Pt(6 if lvl == 0 else 4)
                p.line_spacing = 1.05
                style_run(p, txt, 19 if lvl == 0 else 17, color=INK)
        elif kind == "code":
            for line in payload:
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                p.space_before = p.space_after = Pt(1)
                style_run(p, line or " ", 15, color=CODEINK, mono=True)
    return first


def title_bar(slide, title, idx):
    tb = slide.shapes.add_textbox(Inches(MX), Inches(TITLE_Y), Inches(CW), Inches(TITLE_H))
    tb.text_frame.word_wrap = True
    style_run(tb.text_frame.paragraphs[0], title, 30, True, NAVY)
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(MX), Inches(RULE_Y),
                                    Inches(MX + CW), Inches(RULE_Y))
    ln.line.color.rgb = TEAL; ln.line.width = Pt(2)
    fl = slide.shapes.add_textbox(Inches(MX), Inches(FOOT_Y), Inches(CW), Inches(0.3))
    style_run(fl.text_frame.paragraphs[0], "PyStorm-LCS", 10, color=MUTE)
    pn = slide.shapes.add_textbox(Inches(MX), Inches(FOOT_Y), Inches(CW), Inches(0.3))
    pp = pn.text_frame.paragraphs[0]; pp.alignment = PP_ALIGN.RIGHT
    style_run(pp, str(idx), 10, color=MUTE)


def emu_in(v):
    return v / 914400.0


for n_idx, sl in enumerate(slides, 1):
    body = parse(sl.split("\n"))
    kinds = {k for k, _ in body}
    s = prs.slides.add_slide(BLANK)

    # Title slide -------------------------------------------------------------
    if "big" in kinds:
        y = 2.5
        for k, v in body:
            if k == "big":
                tb = s.shapes.add_textbox(Inches(1), Inches(y), Inches(11.33), Inches(1.3))
                p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                style_run(p, v, 44, True, NAVY); y += 1.3
                ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.2), Inches(y),
                                            Inches(8.13), Inches(y))
                ln.line.color.rgb = TEAL; ln.line.width = Pt(2.5); y += 0.28
            elif k in ("h3", "para"):
                tb = s.shapes.add_textbox(Inches(1), Inches(y), Inches(11.33), Inches(0.75))
                p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                tb.text_frame.word_wrap = True
                style_run(p, v, 25 if k == "h3" else 18, k == "h3",
                          INK if k == "h3" else MUTE)
                y += 0.68 if k == "h3" else 0.55
        continue

    title = next((v for k, v in body if k == "title"), "")
    rest = [(k, v) for k, v in body if k != "title"]
    title_bar(s, title, n_idx)

    # Image slide -------------------------------------------------------------
    if "image" in kinds:
        img = next(v for k, v in rest if k == "image")
        cap = " ".join(v for k, v in rest if k == "para")
        pic = s.shapes.add_picture(str(Path(SRC).parent / img), Inches(0), Inches(BODY_Y),
                                   width=Inches(11.0))
        max_h = 4.9 if cap else 5.7
        if emu_in(pic.height) > max_h:
            w, h = pic.width, pic.height
            pic.height = Inches(max_h); pic.width = Inches(max_h * w / h)
        pic.left = Inches((13.333 - emu_in(pic.width)) / 2)
        if cap:
            cy = emu_in(pic.top + pic.height) + 0.18
            tb = s.shapes.add_textbox(Inches(MX), Inches(cy), Inches(CW), Inches(0.9))
            tb.text_frame.word_wrap = True
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            style_run(p, cap, 17, color=MUTE, italic=True)
        continue

    # Table slide: optional intro text, the table, optional trailing text ------
    if "table" in kinds:
        cur = BODY_Y
        ti = next(i for i, (k, _) in enumerate(rest) if k == "table")
        pre, tbl_rows, post = rest[:ti], rest[ti][1], rest[ti + 1:]
        if pre:
            h = 0.34 * sum((len(v) // 95 + 1) if k == "para" else len(v)
                           for k, v in pre) + 0.1
            tb = s.shapes.add_textbox(Inches(MX), Inches(cur), Inches(CW), Inches(h))
            tb.text_frame.word_wrap = True
            add_blocks(tb.text_frame, pre)
            cur += h + 0.12
        rows = [r for r in tbl_rows if not re.match(r"\s*\|[\s:|-]+\|\s*$", r)]
        grid = [[c.strip() for c in r.strip().strip("|").split("|")] for r in rows]
        nr, nc = len(grid), max(len(r) for r in grid)
        rh = 0.47
        gt = s.shapes.add_table(nr, nc, Inches(MX), Inches(cur), Inches(CW),
                                Inches(rh * nr)).table
        for ri, row in enumerate(grid):
            gt.rows[ri].height = Inches(rh)
            for ci in range(nc):
                cell = gt.cell(ri, ci)
                cell.text = inline(row[ci]) if ci < len(row) else ""
                cell.margin_top = cell.margin_bottom = Pt(2)
                for rn in cell.text_frame.paragraphs[0].runs:
                    rn.font.size = Pt(15); rn.font.name = BODYF
                    rn.font.bold = (ri == 0)
                    rn.font.color.rgb = WHITE if ri == 0 else INK
        cur += rh * nr + 0.15
        if post:
            tb = s.shapes.add_textbox(Inches(MX), Inches(cur), Inches(CW),
                                      Inches(max(0.4, BODY_BOTTOM - cur)))
            tb.text_frame.word_wrap = True
            add_blocks(tb.text_frame, post)
        continue

    # Pure-text slide: one frame spanning the body region ----------------------
    tb = s.shapes.add_textbox(Inches(MX), Inches(BODY_Y), Inches(CW),
                              Inches(BODY_BOTTOM - BODY_Y))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.NONE
    add_blocks(tf, rest)

prs.save(OUT)
print(f"wrote {OUT} with {len(prs.slides._sldIdLst)} slides")
